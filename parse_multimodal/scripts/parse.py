#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
parse_multimodal - 多模态文件解析器
支持的文件类型：PDF、PPTX、DOCX、图片、文本文件
- 文本文件：直接提取文本内容
- 图片文件：调用多模态 API 进行解析
"""

import os
import sys

if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")

import json
import base64
from pathlib import Path
from typing import Dict, Optional, Any, List

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))

try:
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

from common import analyze_image, get_api_config


def parse_pdf(file_path: str) -> Dict:
    """解析 PDF 文件"""
    if not PDF_AVAILABLE:
        return {"status": "error", "error": "pdfplumber not installed. Run: pip install pdfplumber"}

    try:
        text_content = []
        tables = []
        images = []
        
        with pdfplumber.open(file_path) as pdf:
            page_count = len(pdf.pages)
            
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    text_content.append(f"[Page {page_num}]\n{text}")
                
                page_tables = page.extract_tables()
                for table in page_tables:
                    if table:
                        tables.append({"page": page_num, "data": table})
                
                for img in page.images:
                    images.append({
                        "page": page_num,
                        "width": img.get("width"),
                        "height": img.get("height")
                    })

        return {
            "status": "success",
            "file_path": file_path,
            "file_type": "pdf",
            "content": {
                "text": "\n\n".join(text_content),
                "tables": tables,
                "images": [{"page": img["page"], "description": f"[图片: {img['width']:.0f}x{img['height']:.0f}]"} for img in images],
                "metadata": {"page_count": page_count, "table_count": len(tables), "image_count": len(images)}
            }
        }

    except Exception as e:
        return {"status": "error", "error": f"Failed to parse PDF: {str(e)}"}


def parse_pptx(file_path: str) -> Dict:
    """解析 PPTX 文件"""
    if not PPTX_AVAILABLE:
        return {"status": "error", "error": "python-pptx not installed. Run: pip install python-pptx"}

    try:
        prs = Presentation(file_path)
        text_content = []
        images = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
                if hasattr(shape, "image"):
                    images.append({"page": slide_num, "description": "[图片]"})

            if slide_text:
                text_content.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))

        return {
            "status": "success",
            "file_path": file_path,
            "file_type": "presentation",
            "content": {
                "text": "\n\n".join(text_content),
                "images": images,
                "tables": [],
                "metadata": {"slide_count": len(prs.slides), "image_count": len(images)}
            }
        }

    except Exception as e:
        return {"status": "error", "error": f"Failed to parse PPTX: {str(e)}"}


def parse_docx(file_path: str) -> Dict:
    """解析 DOCX 文件"""
    if not DOCX_AVAILABLE:
        return {"status": "error", "error": "python-docx not installed. Run: pip install python-docx"}

    try:
        doc = Document(file_path)
        text_content = []
        tables = []
        images = []

        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text)

        for table in doc.tables:
            table_data = [[cell.text for cell in row.cells] for row in table.rows]
            tables.append({"data": table_data})

        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                images.append({"description": "[图片]"})

        return {
            "status": "success",
            "file_path": file_path,
            "file_type": "document",
            "content": {
                "text": "\n".join(text_content),
                "tables": tables,
                "images": images,
                "metadata": {"paragraph_count": len(doc.paragraphs), "table_count": len(tables), "image_count": len(images)}
            }
        }

    except Exception as e:
        return {"status": "error", "error": f"Failed to parse DOCX: {str(e)}"}


def parse_image(file_path: str, use_vision_api: bool = True) -> Dict:
    """
    解析图片文件
    
    Args:
        file_path: 图片路径
        use_vision_api: 是否使用多模态 API 解析图片内容
        
    Returns:
        解析结果
    """
    if not PIL_AVAILABLE:
        return {"status": "error", "error": "Pillow not installed. Run: pip install Pillow"}

    try:
        img = Image.open(file_path)
        
        result = {
            "status": "success",
            "file_path": file_path,
            "file_type": "image",
            "content": {
                "text": "",
                "images": [],
                "tables": [],
                "metadata": {
                    "width": img.width,
                    "height": img.height,
                    "format": img.format,
                    "mode": img.mode
                }
            }
        }
        
        if use_vision_api:
            vision_prompt = """请详细分析这张教学材料图片，提取以下信息：

1. **文字内容**：识别并提取图片中的所有文字
2. **图表描述**：如果有图表，请详细描述
3. **公式识别**：如果有公式，请准确识别
4. **表格内容**：如果有表格，请提取数据
5. **整体说明**：简要说明图片的教学主题

请用中文回答。"""

            vision_result = analyze_image(image_path=file_path, prompt=vision_prompt)
            
            if vision_result.get("status") == "success":
                description = vision_result.get("content", "")
                result["content"]["text"] = description
                result["content"]["images"].append({
                    "page": 1,
                    "description": description,
                    "metadata": {"width": img.width, "height": img.height, "format": img.format}
                })
            else:
                result["content"]["text"] = f"[图片文件: {Path(file_path).name}]"
                result["content"]["images"].append({
                    "page": 1,
                    "description": f"[图片解析失败: {vision_result.get('error')}]",
                    "metadata": {"width": img.width, "height": img.height}
                })
        else:
            result["content"]["text"] = f"[图片文件: {Path(file_path).name}, 尺寸: {img.width}x{img.height}]"
            result["content"]["images"].append({
                "page": 1,
                "description": f"[图片: {img.width}x{img.height}, {img.format}]",
                "metadata": {"width": img.width, "height": img.height, "format": img.format}
            })

        return result

    except Exception as e:
        return {"status": "error", "error": f"Failed to parse image: {str(e)}"}


def parse_text_file(file_path: str) -> Dict:
    """解析文本文件"""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
        
        ext = Path(file_path).suffix.lower()
        file_type = "markdown" if ext == ".md" else "text"
        
        return {
            "status": "success",
            "file_path": file_path,
            "file_type": file_type,
            "content": {
                "text": text,
                "images": [],
                "tables": [],
                "metadata": {"line_count": len(text.split("\n")), "char_count": len(text)}
            }
        }
    except UnicodeDecodeError:
        try:
            with open(file_path, "r", encoding="gbk") as f:
                text = f.read()
            return {
                "status": "success",
                "file_path": file_path,
                "file_type": "text",
                "content": {"text": text, "images": [], "tables": [], "metadata": {"encoding": "gbk"}}
            }
        except Exception as e:
            return {"status": "error", "error": f"Failed to read text file: {str(e)}"}
    except Exception as e:
        return {"status": "error", "error": f"Failed to parse text file: {str(e)}"}


def get_file_type(file_path: str) -> str:
    """根据文件扩展名判断文件类型"""
    ext = Path(file_path).suffix.lower()
    type_map = {
        ".pdf": "pdf", ".pptx": "pptx", ".ppt": "ppt_legacy",
        ".docx": "docx", ".doc": "doc_legacy",
        ".xlsx": "xlsx", ".xls": "xls_legacy",
        ".png": "image", ".jpg": "image", ".jpeg": "image",
        ".gif": "image", ".bmp": "image", ".webp": "image",
        ".txt": "text", ".md": "text", ".markdown": "text",
    }
    return type_map.get(ext, "unknown")


def parse_file(file_path: str, use_vision_api: bool = True) -> Dict:
    """
    主解析函数
    
    Args:
        file_path: 文件路径
        use_vision_api: 是否使用多模态 API 解析图片（仅对图片文件有效）
        
    Returns:
        解析结果
    """
    if not os.path.exists(file_path):
        return {"status": "error", "error": "File not found"}

    if not os.path.isfile(file_path):
        return {"status": "error", "error": "Path is not a file"}

    file_type = get_file_type(file_path)

    parsers = {
        "pdf": lambda: parse_pdf(file_path),
        "pptx": lambda: parse_pptx(file_path),
        "docx": lambda: parse_docx(file_path),
        "image": lambda: parse_image(file_path, use_vision_api),
        "text": lambda: parse_text_file(file_path),
    }

    legacy_formats = {
        "ppt_legacy": "旧版 .PPT 格式不支持，请将文件另存为 .pptx 格式",
        "doc_legacy": "旧版 .DOC 格式不支持，请将文件另存为 .docx 格式",
        "xlsx": "Excel 文件暂不支持解析",
        "xls_legacy": "旧版 .XLS 格式不支持，请将文件另存为 .xlsx 格式",
    }

    if file_type in legacy_formats:
        return {"status": "error", "error": legacy_formats[file_type], "file_path": file_path, "file_type": file_type}

    if file_type in parsers:
        return parsers[file_type]()

    return {"status": "error", "error": f"Unsupported file type: {file_type}"}


def main():
    """主函数"""
    import argparse
    parser = argparse.ArgumentParser(description="多模态文件解析")
    parser.add_argument("file_path", help="要解析的文件路径")
    parser.add_argument("--no-vision", action="store_true", help="禁用图片的多模态 API 解析")
    
    args = parser.parse_args()
    result = parse_file(args.file_path, use_vision_api=not args.no_vision)
    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
